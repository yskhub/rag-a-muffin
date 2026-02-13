"""
RAG-a-Muffin: Stunning PowerPoint Presentation Generator
Creates a professional, dark-themed enterprise presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───
BG_DARK = RGBColor(0x02, 0x06, 0x17)       # #020617
BG_CARD = RGBColor(0x0F, 0x17, 0x2A)       # #0f172a
PRIMARY = RGBColor(0x00, 0xF0, 0xFF)        # #00f0ff (cyan)
ACCENT = RGBColor(0x0F, 0xFF, 0xC1)         # #0fffc1 (green)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)    # slate-200
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)      # slate-400
DIM_GRAY = RGBColor(0x64, 0x74, 0x8B)      # slate-500
EMERALD = RGBColor(0x34, 0xD3, 0x99)       # emerald-400
AMBER = RGBColor(0xFB, 0xBF, 0x24)         # amber-400
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)        # violet-400
RED = RGBColor(0xF8, 0x71, 0x71)           # red-400

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helper Functions ───

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=16, color=LIGHT_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p

def add_shape_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    return shape

def add_dot(slide, left, top, color=EMERALD, size=0.15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top),
        Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_slide_number(slide, num, total):
    add_text_box(slide, 12.0, 7.0, 1.0, 0.4, f"{num}/{total}",
                 font_size=10, color=DIM_GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 12

# ═══════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide)

# Accent line at top
add_accent_line(slide, 0, 0, 13.333)

# Diamond icon
add_text_box(slide, 5.5, 1.8, 2.3, 1.0, "◆", font_size=60, color=PRIMARY,
             alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide, 1.5, 2.8, 10.3, 1.2, "RAG-a-Muffin", font_size=54,
             color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, 2.0, 3.9, 9.3, 0.8,
             "AI Context Engine — Retrieval-Augmented Generation",
             font_size=22, color=PRIMARY, alignment=PP_ALIGN.CENTER)

# Description
add_text_box(slide, 2.5, 4.8, 8.3, 1.0,
             "An enterprise-grade AI assistant powered by RAG,\nbuilt entirely on free-tier services. Total cost: $0.00",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom info
add_shape_card(slide, 3.5, 6.0, 6.3, 0.8, BG_CARD)
add_text_box(slide, 3.7, 6.1, 5.9, 0.6,
             "🌐 rag-a-muffin.vercel.app    |    📅 February 2026",
             font_size=14, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "The Problem", font_size=36,
             color=WHITE, bold=True)

# Problem cards
problems = [
    ("❌", "Static Knowledge", "Traditional AI chatbots only know what they were trained on.\nThey can't answer questions about YOUR specific data."),
    ("❌", "Hallucination Risk", "Without grounding in real data, AI confidently makes up\ninformation that sounds correct but is completely wrong."),
    ("❌", "Cost Barrier", "Enterprise AI solutions cost hundreds per month.\nSmall teams and students can't afford them."),
]

for i, (icon, title, desc) in enumerate(problems):
    y = 2.0 + i * 1.7
    add_shape_card(slide, 0.8, y, 11.7, 1.4)
    add_text_box(slide, 1.2, y + 0.15, 0.6, 0.5, icon, font_size=24,
                 color=RED)
    add_text_box(slide, 2.0, y + 0.1, 4.0, 0.5, title, font_size=20,
                 color=WHITE, bold=True)
    add_text_box(slide, 2.0, y + 0.6, 10.0, 0.7, desc, font_size=14,
                 color=MID_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 3: The Solution — RAG
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "The Solution: RAG", font_size=36,
             color=WHITE, bold=True)
add_text_box(slide, 0.8, 1.5, 10, 0.5,
             "Retrieval-Augmented Generation — AI grounded in YOUR data",
             font_size=16, color=PRIMARY)

steps = [
    ("1", "Upload", "Documents are chunked and\nconverted to vector embeddings", PRIMARY),
    ("2", "Retrieve", "User queries are matched against\nthe vector database for relevant context", EMERALD),
    ("3", "Augment", "Retrieved context is injected into\nthe AI prompt as grounding data", AMBER),
    ("4", "Generate", "Gemini generates an answer based\non YOUR data, not hallucination", VIOLET),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.8 + i * 3.1
    add_shape_card(slide, x, 2.5, 2.8, 3.5)
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 1.05), Inches(2.8),
        Inches(0.7), Inches(0.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x05, 0x10, 0x1E)
    circle.line.color.rgb = color
    circle.line.width = Pt(2)
    
    add_text_box(slide, x + 1.05, 2.85, 0.7, 0.6, num, font_size=24,
                 color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, 3.8, 2.2, 0.5, title, font_size=20,
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 4.4, 2.4, 1.0, desc, font_size=13,
                 color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow connectors between steps
for i in range(3):
    x = 3.7 + i * 3.1
    add_text_box(slide, x, 3.8, 0.5, 0.5, "→", font_size=28,
                 color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom callout
add_shape_card(slide, 2.5, 6.3, 8.3, 0.7)
add_text_box(slide, 2.7, 6.4, 7.9, 0.5,
             "💡 The AI doesn't guess — it retrieves. Every answer is grounded in real documents.",
             font_size=14, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 4: Architecture
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "System Architecture", font_size=36,
             color=WHITE, bold=True)

# Architecture layers
layers = [
    ("Frontend", "React + Zustand + Vite", "Vercel (Free)", PRIMARY, 0.8),
    ("API Layer", "FastAPI + Python", "Render (Free)", EMERALD, 3.2),
    ("AI Engine", "Google Gemini 2.0 Flash", "Free — 15 RPM", AMBER, 5.6),
    ("Vector DB", "ChromaDB", "In-Memory (Free)", VIOLET, 8.0),
]

for name, tech, hosting, color, x in layers:
    add_shape_card(slide, x, 2.2, 3.8 if x < 7 else 4.5, 2.0)
    add_dot(slide, x + 0.3, 2.45, color)
    add_text_box(slide, x + 0.6, 2.3, 3.0, 0.4, name, font_size=18,
                 color=WHITE, bold=True)
    add_text_box(slide, x + 0.3, 2.9, 3.2, 0.4, tech, font_size=14,
                 color=color)
    add_text_box(slide, x + 0.3, 3.4, 3.2, 0.4, hosting, font_size=12,
                 color=DIM_GRAY)

# Data flow arrows
for x_pos in [4.7, 7.1, 9.7 ]:
    if x_pos < 10:
        add_text_box(slide, x_pos - 0.1, 2.8, 0.8, 0.5, "⟶", font_size=24,
                     color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

# Cost box
add_shape_card(slide, 3.5, 5.0, 6.3, 1.8)
add_text_box(slide, 3.8, 5.2, 5.7, 0.5, "Total Infrastructure Cost",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 3.8, 5.7, 5.7, 0.8, "$0.00 / month",
             font_size=42, color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 3.8, 6.3, 5.7, 0.4, "No credit card required. 100% free-tier services.",
             font_size=13, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 4, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 5: Tech Stack
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "Technology Stack", font_size=36,
             color=WHITE, bold=True)

stack = [
    ("Frontend", [
        ("React 18", "Component-based UI framework"),
        ("Zustand", "Lightweight state management with persistence"),
        ("Vite", "Lightning-fast build tool"),
        ("React-Markdown", "Markdown rendering for AI responses"),
    ], PRIMARY),
    ("Backend", [
        ("FastAPI", "High-performance Python API framework"),
        ("Google Gemini", "Free LLM — 2.0 Flash (latest)"),
        ("ChromaDB", "Open-source vector embedding database"),
        ("In-Memory Store", "Session-based conversation memory"),
    ], EMERALD),
    ("DevOps", [
        ("Vercel", "Frontend hosting with auto-deploy"),
        ("Render", "Backend hosting with free tier"),
        ("GitHub", "Version control & CI/CD trigger"),
        ("CORS + GZip", "Security & performance middleware"),
    ], VIOLET),
]

for col, (category, items, color) in enumerate(stack):
    x = 0.8 + col * 4.1
    add_text_box(slide, x, 2.0, 3.6, 0.5, category, font_size=22,
                 color=color, bold=True)
    
    for j, (tech, desc) in enumerate(items):
        y = 2.7 + j * 1.1
        add_shape_card(slide, x, y, 3.8, 0.9)
        add_dot(slide, x + 0.25, y + 0.35, color)
        add_text_box(slide, x + 0.55, y + 0.1, 3.0, 0.3, tech,
                     font_size=15, color=WHITE, bold=True)
        add_text_box(slide, x + 0.55, y + 0.45, 3.0, 0.4, desc,
                     font_size=11, color=DIM_GRAY)

add_slide_number(slide, 5, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 6: RAG Pipeline Flow
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "RAG Pipeline — How It Works",
             font_size=36, color=WHITE, bold=True)

pipeline_steps = [
    ("Step 1", "User Query", "User sends a question\nthrough the chat interface", "💬", PRIMARY),
    ("Step 2", "Vector Search", "ChromaDB searches for\nrelevant document chunks", "🔍", EMERALD),
    ("Step 3", "Context Building", "Top-5 relevant fragments\nare assembled as context", "📋", AMBER),
    ("Step 4", "Memory Injection", "Last 5 conversation messages\nare added for continuity", "🧠", VIOLET),
    ("Step 5", "LLM Generation", "Gemini generates answer\ngrounded in retrieved context", "⚡", PRIMARY),
    ("Step 6", "Source Citation", "Response includes source\nbadges proving data origin", "📎", EMERALD),
]

for i, (step, title, desc, icon, color) in enumerate(pipeline_steps):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 2.2 + row * 2.6

    add_shape_card(slide, x, y, 3.8, 2.2)
    add_text_box(slide, x + 0.3, y + 0.2, 0.5, 0.4, icon, font_size=22)
    add_text_box(slide, x + 0.9, y + 0.15, 2.6, 0.3, step, font_size=11,
                 color=color, bold=True)
    add_text_box(slide, x + 0.3, y + 0.7, 3.2, 0.4, title, font_size=18,
                 color=WHITE, bold=True)
    add_text_box(slide, x + 0.3, y + 1.2, 3.2, 0.8, desc, font_size=13,
                 color=MID_GRAY)

add_slide_number(slide, 6, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 7: Key Features
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "Key Features", font_size=36,
             color=WHITE, bold=True)

features = [
    ("🎯", "Quick Actions", "Pre-built prompt suggestions guide users", PRIMARY),
    ("⏱️", "Response Timer", "Tracks generation time per response", EMERALD),
    ("🟢", "Live Health Check", "Real-time backend monitoring every 30s", AMBER),
    ("📝", "Markdown Rendering", "AI responses with bold, lists, code blocks", VIOLET),
    ("👍", "Feedback System", "Thumbs up/down on every AI response", PRIMARY),
    ("📊", "Real Metrics", "Live document count from backend API", EMERALD),
    ("📋", "Chat Export", "Copy entire conversation to clipboard", AMBER),
    ("⌨️", "Keyboard Shortcuts", "Ctrl+K clear, Ctrl+E export, Enter send", VIOLET),
    ("💬", "Session History", "Auto-saved sessions with load/delete", PRIMARY),
]

for i, (icon, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 2.0 + row * 1.6

    add_shape_card(slide, x, y, 3.8, 1.3)
    add_text_box(slide, x + 0.3, y + 0.15, 0.5, 0.4, icon, font_size=20)
    add_text_box(slide, x + 0.9, y + 0.15, 2.6, 0.4, title, font_size=16,
                 color=WHITE, bold=True)
    add_text_box(slide, x + 0.3, y + 0.65, 3.2, 0.5, desc, font_size=12,
                 color=MID_GRAY)

add_slide_number(slide, 7, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 8: UI Design
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "Enterprise UI Design", font_size=36,
             color=WHITE, bold=True)
add_text_box(slide, 0.8, 1.5, 10, 0.5,
             "Professional dark theme with glassmorphism, live metrics, and responsive layout",
             font_size=16, color=MID_GRAY)

ui_elements = [
    ("3-Column Layout", "Chat (7) + Metrics (3) + AI Core (2)\n12-column responsive grid", PRIMARY),
    ("Glassmorphism Cards", "Semi-transparent cards with backdrop blur\nand subtle borders for depth", EMERALD),
    ("Live Dashboard", "Real-time metrics, subsystem status,\nand activity log updating continuously", AMBER),
    ("Chat Bubbles", "User (right, cyan) and AI (left, slate)\nwith inline source citations", VIOLET),
    ("Typography System", "Space Grotesk (headings), Inter (body),\nJetBrains Mono (code)", PRIMARY),
    ("Micro-Animations", "Float, slide-up, pulse-glow, scanline\nfor premium interactive feel", EMERALD),
]

for i, (title, desc, color) in enumerate(ui_elements):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.3
    y = 2.3 + row * 1.7

    add_shape_card(slide, x, y, 5.8, 1.4)
    add_dot(slide, x + 0.3, y + 0.25, color, 0.12)
    add_text_box(slide, x + 0.6, y + 0.1, 5.0, 0.4, title, font_size=17,
                 color=WHITE, bold=True)
    add_text_box(slide, x + 0.6, y + 0.6, 5.0, 0.7, desc, font_size=13,
                 color=MID_GRAY)

add_slide_number(slide, 8, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 9: Live Demo Plan
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "Live Demo", font_size=36,
             color=WHITE, bold=True)
add_text_box(slide, 0.8, 1.5, 10, 0.5,
             "rag-a-muffin.vercel.app",
             font_size=20, color=PRIMARY)

demo_steps = [
    ("1", "Seed Knowledge Base", "Admin → Initialize Sample Data\n(Products, Policies, Shipping)", "30s"),
    ("2", "Product Query", '"What headphones do you sell?"\n→ Retrieves from Catalog source', "45s"),
    ("3", "Policy Query", '"What is your return policy?"\n→ Retrieves from FAQ source', "30s"),
    ("4", "Memory Test", '"Can I return those headphones?"\n→ Combines both contexts', "30s"),
    ("5", "Honesty Test", '"Do you sell laptops?"\n→ AI says "I don\'t have that info"', "30s"),
    ("6", "Features", "Feedback, Export, History,\nKeyboard Shortcuts demo", "30s"),
]

for i, (num, title, desc, duration) in enumerate(demo_steps):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.3
    y = 2.3 + row * 1.7

    add_shape_card(slide, x, y, 5.8, 1.4)
    # Step number
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.2), Inches(y + 0.3),
        Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x00, 0x20, 0x30)
    circle.line.color.rgb = PRIMARY
    circle.line.width = Pt(1.5)
    add_text_box(slide, x + 0.2, y + 0.32, 0.5, 0.45, num, font_size=18,
                 color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + 0.9, y + 0.1, 4.0, 0.4, title, font_size=16,
                 color=WHITE, bold=True)
    add_text_box(slide, x + 0.9, y + 0.55, 4.0, 0.7, desc, font_size=12,
                 color=MID_GRAY)
    add_text_box(slide, x + 5.0, y + 0.2, 0.6, 0.3, duration, font_size=11,
                 color=DIM_GRAY, alignment=PP_ALIGN.RIGHT)

add_slide_number(slide, 9, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 10: Challenges & Solutions
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "Challenges & Solutions", font_size=36,
             color=WHITE, bold=True)

challenges = [
    ("Rate Limiting", "Gemini free tier: 15 RPM",
     "Self-limit to 10 RPM + 4s min gap\n+ exponential backoff (15→30→60s)", RED, EMERALD),
    ("Cold Starts", "Render free tier sleeps after 15 min",
     "120s frontend timeout + live health\ncheck with visual status indicator", AMBER, EMERALD),
    ("Stale Errors", "Rate-limit errors cached in localStorage",
     "Transient error filter on rehydration\n+ manual Clear Chat button", RED, EMERALD),
    ("CSS in Production", "Tailwind 4 not processing classes",
     "@tailwindcss/vite plugin added to\nvite.config.js — fixed all styling", AMBER, EMERALD),
]

for i, (challenge, problem, solution, prob_color, sol_color) in enumerate(challenges):
    y = 2.0 + i * 1.3
    
    add_shape_card(slide, 0.8, y, 11.7, 1.1)
    add_text_box(slide, 1.2, y + 0.1, 2.5, 0.35, challenge, font_size=16,
                 color=WHITE, bold=True)
    
    add_text_box(slide, 1.2, y + 0.55, 4.0, 0.5, f"Problem: {problem}",
                 font_size=12, color=prob_color)
    
    add_text_box(slide, 6.5, y + 0.15, 0.5, 0.4, "→", font_size=20,
                 color=DIM_GRAY, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, 7.2, y + 0.15, 5.0, 0.8, f"✅ {solution}",
                 font_size=12, color=sol_color)

add_slide_number(slide, 10, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 11: Results & Impact
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0.8, 0.6, 2.0)

add_text_box(slide, 0.8, 0.8, 10, 0.8, "Results & Impact", font_size=36,
             color=WHITE, bold=True)

metrics_data = [
    ("$0.00", "Monthly Cost", "100% free-tier\ninfrastructure", EMERALD),
    ("< 3s", "Response Time", "Average query-to-answer\nwith RAG pipeline", PRIMARY),
    ("9", "UI Features", "Production-grade\nenterprise features", AMBER),
    ("0", "Hallucinations", "RAG-grounded responses\nwith source citations", VIOLET),
]

for i, (value, label, desc, color) in enumerate(metrics_data):
    x = 0.8 + i * 3.1
    add_shape_card(slide, x, 2.2, 2.8, 2.5)
    add_text_box(slide, x + 0.2, 2.5, 2.4, 0.8, value, font_size=36,
                 color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 3.3, 2.4, 0.4, label, font_size=16,
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 3.8, 2.4, 0.6, desc, font_size=12,
                 color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

# Key achievement
add_shape_card(slide, 0.8, 5.2, 11.7, 1.8)
add_text_box(slide, 1.2, 5.4, 11.0, 0.5, "Key Achievement", font_size=18,
             color=PRIMARY, bold=True)

achievements = [
    "✅  Built a fully functional RAG system with document ingestion, vector search, and contextual AI generation",
    "✅  Enterprise-grade UI with live monitoring, session management, markdown rendering, and feedback system",
    "✅  Zero-cost deployment using Vercel + Render + Gemini free tiers — no credit card required",
    "✅  Production patterns: rate limiting, error recovery, exponential backoff, state persistence",
]

txBox = add_text_box(slide, 1.2, 5.9, 11.0, 0.3, achievements[0], font_size=13, color=LIGHT_GRAY)
for ach in achievements[1:]:
    add_para(txBox.text_frame, ach, font_size=13, color=LIGHT_GRAY, space_before=Pt(4))

add_slide_number(slide, 11, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SLIDE 12: Thank You / Q&A
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, 0, 0, 13.333)

add_text_box(slide, 1.5, 1.5, 10.3, 1.0, "Thank You", font_size=54,
             color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2.0, 2.8, 9.3, 0.6,
             "Questions & Discussion",
             font_size=24, color=PRIMARY, alignment=PP_ALIGN.CENTER)

# Links card
add_shape_card(slide, 3.0, 3.8, 7.3, 2.5)

links = [
    ("🌐  Live App", "rag-a-muffin.vercel.app"),
    ("⚙️  Backend API", "rag-a-muffin.onrender.com"),
    ("📂  Source Code", "github.com/yskhub/RAG-a-Muffin"),
    ("🛠️  Tech Stack", "React · FastAPI · Gemini · ChromaDB"),
]

for i, (label, value) in enumerate(links):
    y = 4.0 + i * 0.5
    add_text_box(slide, 3.5, y, 2.5, 0.4, label, font_size=14,
                 color=MID_GRAY)
    add_text_box(slide, 6.2, y, 3.8, 0.4, value, font_size=14,
                 color=LIGHT_GRAY, bold=True)

# Bottom tagline
add_text_box(slide, 2.0, 6.5, 9.3, 0.5,
             "Built entirely on free-tier services  ·  Total cost: $0.00",
             font_size=14, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 12, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "RAG-a-Muffin_Presentation.pptx")
prs.save(output_path)
print(f"\n✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {TOTAL_SLIDES}")
print(f"🎨 Theme: Dark enterprise with cyan accent")
